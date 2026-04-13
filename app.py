import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from PIL import Image
import io
import datetime

# --- DÖKÜMAN VE VERİ İŞLEME KÜTÜPHANELERİ ---
from docx import Document # Word için
from pptx import Presentation # PowerPoint için
import pandas as pd # Excel için
from fpdf import FPDF # PDF için
from docx2python import docx2python # Word okuma için
from pptx import Presentation as ReadPPTX # PPTX okuma için

# --- SİSTEM AYARLARI ---
st.set_page_config(page_title="Eren AI | Akademik Portal", page_icon="🛡️", layout="wide")

if "GOOGLE_API_KEY" in st.secrets:
    genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])
else:
    st.error("API Anahtarı bulunamadı! Lütfen secrets.toml dosyasını kontrol edin.")
    st.stop()

model = genai.GenerativeModel('gemini-3-flash-preview')

# --- ÖZEL EREN FEN VE TEKNOLOJİ LİSESİ BİLGİ TABANI ---
OKUL_BILGILERI = "Kurum: Özel Eren Fen ve Teknoloji Lisesi | Web: https://eren.k12.tr/"

# --- ANAYASAL DÖKÜMAN ÜRETİM MOTORLARI (2. ETAP) ---

def create_word_file(text, topic):
    doc = Document()
    doc.add_heading('Özel Eren Fen ve Teknoloji Lisesi', 0)
    doc.add_heading(f'Akademik Çalışma Fasikülü: {topic}', level=1)
    doc.add_paragraph(f"Tarih: {datetime.datetime.now().strftime('%d/%m/%Y')}")
    doc.add_paragraph(text)
    doc.add_paragraph("\n\n© Eren AI Education - Eren Yapay Zeka Sistemleri")
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

def create_pdf_file(text, topic):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_font('Arial', '', '', uni=True) 
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Özel Eren Fen ve Teknoloji Lisesi", ln=True, align='C')
    pdf.ln(10)
    pdf.multi_cell(0, 10, txt=text.encode('latin-1', 'ignore').decode('latin-1'))
    return pdf.output(dest='S').encode('latin-1', 'ignore')

def create_pptx_file(text, topic):
    prs = Presentation()
    # Kapak Slaytı
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = topic
    slide.placeholders[1].text = "Eren AI Education Akademik Sunumu"
    
    # İçerik Slaytı (Özet)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Akademik Analiz Özeti"
    slide.placeholders[1].text = text[:1000] + "..."
    
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()

def create_excel_file(text):
    # Metni satırlara bölüp basit bir tablo oluşturur
    data = {"Akademik Analiz İçeriği": [text[:500]], "Oluşturulma": [datetime.datetime.now()]}
    df = pd.DataFrame(data)
    bio = io.BytesIO()
    with pd.ExcelWriter(bio, engine='openpyxl') as writer:
        df.to_excel(writer, index=False)
    return bio.getvalue()

# --- DOSYA OKUMA MOTORLARI ---

def read_file_content(uploaded_file):
    content = ""
    if uploaded_file.type == "application/pdf":
        reader = PdfReader(uploaded_file)
        for page in reader.pages:
            content += page.extract_text() + "\n"
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        content = docx2python(uploaded_file).text
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = ReadPPTX(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
    elif uploaded_file.type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "text/csv"]:
        df = pd.read_excel(uploaded_file) if uploaded_file.name.endswith('xlsx') else pd.read_csv(uploaded_file)
        content = df.to_string()
    return content

# --- SOL MENÜ (ANAYASAL VE SABİT) ---
def sidebar_ciz():
    with st.sidebar:
        try:
            st.image("Logo.png", use_container_width=True)
        except:
            st.subheader("🛡️ Eren AI")
        
        st.markdown("---")
        st.markdown("### **🛡️ Eren AI Education**")
        st.markdown("**Eren Yapay Zeka Sistemleri**")
        st.success("**Anayasal Mod:** Her soru için bireysel ve derin analiz zorunludur.")
        
        st.markdown("### **📖 Sistem Rehberi**")
        st.info("""
        1. **Akademik dökümanlarınızı veya ödev dosyalarınızı** "Dosya Yükleme" alanından sisteme güvenle iletebilirsiniz.
        2. Analiz edilmesini istediğiniz tüm hususları **metin alanına girerek** Eren AI ile akademik etkileşim başlatabilirsiniz.
        """)
        
        st.divider()
        st.caption("© 2026 Eren Eğitim Kurumları")

sidebar_ciz()

# --- SOHBET GEÇMİŞİ VE ARAYÜZ ---
if "messages" not in st.session_state:
    st.session_state.messages = []

chat_area = st.container()
with chat_area:
    for m in st.session_state.messages:
        with st.chat_message(m["role"]):
            st.markdown(m["content"])

# --- GİRİŞ PANELİ ---
with st.container():
    st.write("---")
    dosya = st.file_uploader("Dosya Yükleme", type=['pdf','docx','xlsx','pptx','csv','png','jpg','jpeg'], 
                             key="uploader_main", 
                             label_visibility="collapsed")
    
    soru = st.chat_input("Eren AI'a sormak istediğin soruyu bu alana girebilirsiniz.")

# --- ANA İŞLEMCİ ---
if soru:
    st.session_state.messages.append({"role": "user", "content": soru})
    with chat_area:
        with st.chat_message("user"):
            st.markdown(soru)

    with chat_area:
        with st.chat_message("assistant"):
            durum = st.status("🛡️ Eren AI Akademik Analiz ve Döküman Motoru Çalışıyor...")
            
            try:
                # ANAYASAL TALİMAT (1. VE 2. ETAP BÜTÜNLEŞİK)
                system_instruction = f"""
                Sen Eren AI, Özel Eren Fen ve Teknoloji Lisesi'nin Yapay Zekasısın. {OKUL_BILGILERI}
                Cevaplarına başlarken "Eren AI, Özel Eren Fen ve Teknoloji Lisesi'nin Yapay Zekası olarak size hizmet ediyorum." ifadesini kullan.

                🛡️ AKADEMİK ANAYASA:
                1. MATERYALDEKİ TÜM SORULARI TEKER TEKER ANALİZ ET. Asla toplu cevap verme. Her soru için ## SORU [NO] başlığı aç.
                2. DERİNLİK ZORUNLUDUR. Bilimsel yasaları bir ders notu zenginliğinde anlat.
                3. CEVAP VERME, ÖĞRET. Şıkları analiz et ama doğruyu söyleme, Sokratik yöntemi uygula.
                4. ZORUNLU FORMAT: Her soru için '📚 Kavramsal ve Bilimsel Derinlik', '🔍 Analitik İnceleme ve Çözüm Stratejisi' ve '💡 Analitik Sorgulama ve Sentez' başlıklarını kullan.
                """
                
                input_data = [system_instruction, soru]
                
                if dosya:
                    if dosya.type.startswith("image/"):
                        input_data.append(Image.open(dosya))
                    else:
                        file_text = read_file_content(dosya)
                        input_data.append(f"ANALİZ EDİLECEK DOSYA İÇERİĞİ:\n{file_text}")

                # AI Analizi Başlatır
                full_response = ""
                response = model.generate_content(input_data)
                full_response = response.text
                
                # EKRANA YAZDIRMA (1. ETAP)
                st.markdown(full_response)
                st.session_state.messages.append({"role": "assistant", "content": full_response})
                
                # DÖKÜMAN OLUŞTURMA PANELİ (2. ETAP)
                st.divider()
                st.write("### 📥 Akademik Çıktı Merkezi")
                st.caption("Analiz sonuçlarını kurumsal döküman formatlarında indirin:")
                
                c1, c2, c3, c4 = st.columns(4)
                
                # Word Hazırla
                word_bytes = create_word_file(full_response, "Akademik Analiz Raporu")
                c1.download_button("📄 Word Fasikül", word_bytes, f"ErenAI_Analiz_{datetime.date.today()}.docx")
                
                # PDF Hazırla
                pdf_bytes = create_pdf_file(full_response, "Akademik Analiz Raporu")
                c2.download_button("📕 PDF Rapor", pdf_bytes, f"ErenAI_Analiz_{datetime.date.today()}.pdf")
                
                # PPTX Hazırla
                pptx_bytes = create_pptx_file(full_response, "Akademik Sunum")
                c3.download_button("📽️ PPTX Sunum", pptx_bytes, f"ErenAI_Sunum_{datetime.date.today()}.pptx")
                
                # Excel Hazırla
                excel_bytes = create_excel_file(full_response)
                c4.download_button("📊 Excel Veri", excel_bytes, f"ErenAI_Veri_{datetime.date.today()}.xlsx")
                
                durum.update(label="✅ Analiz ve Dökümanlar Hazır", state="complete")
                
            except Exception as e:
                durum.update(label="❌ Sistem Hatası", state="error")
                st.error(f"Teknik bir sorun oluştu: {str(e)}")
